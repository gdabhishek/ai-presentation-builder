"""
Professional PowerPoint Templates and Design System
This module provides pre-built templates and design utilities for creating
beautiful, professional PowerPoint presentations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, List, Tuple, Optional
import os
from pathlib import Path
import logging

logger = logging.getLogger(__name__)

class PPTDesignSystem:
    """Professional design system for PowerPoint presentations"""
    
    # Color schemes
    COLORS = {
        "business_blue": {
            "primary": RGBColor(0, 120, 212),      # Microsoft Blue
            "secondary": RGBColor(40, 40, 40),     # Dark Gray
            "accent": RGBColor(0, 178, 148),       # Teal
            "background": RGBColor(255, 255, 255), # White
            "text": RGBColor(32, 32, 32)           # Near Black
        },
        "modern_purple": {
            "primary": RGBColor(102, 45, 145),     # Purple
            "secondary": RGBColor(68, 68, 68),     # Gray
            "accent": RGBColor(255, 185, 0),       # Amber
            "background": RGBColor(255, 255, 255), # White
            "text": RGBColor(33, 33, 33)           # Dark Gray
        },
        "tech_green": {
            "primary": RGBColor(46, 125, 50),      # Green
            "secondary": RGBColor(55, 71, 79),     # Blue Gray
            "accent": RGBColor(255, 152, 0),       # Orange
            "background": RGBColor(255, 255, 255), # White
            "text": RGBColor(33, 33, 33)           # Dark Gray
        },
        "corporate_red": {
            "primary": RGBColor(183, 28, 28),      # Red
            "secondary": RGBColor(66, 66, 66),     # Gray
            "accent": RGBColor(0, 188, 212),       # Cyan
            "background": RGBColor(255, 255, 255), # White
            "text": RGBColor(33, 33, 33)           # Dark Gray
        }
    }
    
    # Font settings
    FONTS = {
        "title": {
            "name": "Calibri",
            "size": Pt(44),
            "bold": True
        },
        "subtitle": {
            "name": "Calibri",
            "size": Pt(20),
            "bold": False
        },
        "heading": {
            "name": "Calibri",
            "size": Pt(32),
            "bold": True
        },
        "body": {
            "name": "Calibri",
            "size": Pt(16),
            "bold": False
        },
        "caption": {
            "name": "Calibri",
            "size": Pt(14),
            "bold": False
        }
    }

class PPTTemplate:
    """Base class for PowerPoint templates"""
    
    def __init__(self, color_scheme: str = "business_blue"):
        self.prs = Presentation()
        self.color_scheme = PPTDesignSystem.COLORS.get(color_scheme, PPTDesignSystem.COLORS["business_blue"])
        self.fonts = PPTDesignSystem.FONTS
        
    def apply_font_style(self, text_frame, style_name: str, color: Optional[RGBColor] = None):
        """Apply font styling to a text frame"""
        font_config = self.fonts[style_name]
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = font_config["name"]
            paragraph.font.size = font_config["size"]
            paragraph.font.bold = font_config["bold"]
            if color:
                paragraph.font.color.rgb = color
            else:
                paragraph.font.color.rgb = self.color_scheme["text"]
    
    def create_title_slide(self, title: str, subtitle: str = "") -> None:
        """Create a professional title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        self.apply_font_style(title_shape.text_frame, "title", self.color_scheme["primary"])
        
        # Set subtitle if provided
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            self.apply_font_style(subtitle_shape.text_frame, "subtitle", self.color_scheme["secondary"])
        
        # Add accent line
        self._add_accent_line(slide)
    
    def create_content_slide(self, title: str, content: List[str], layout_type: str = "bullet") -> None:
        """Create a content slide with bullet points or numbered list"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        self.apply_font_style(title_shape.text_frame, "heading", self.color_scheme["primary"])
        
        # Set content
        content_shape = slide.placeholders[1]
        
        if layout_type == "bullet":
            content_text = "\n".join([f"• {item}" for item in content])
        elif layout_type == "numbered":
            content_text = "\n".join([f"{i+1}. {item}" for i, item in enumerate(content)])
        else:
            content_text = "\n".join(content)
        
        content_shape.text = content_text
        self.apply_font_style(content_shape.text_frame, "body")
        
        # Add accent element
        self._add_accent_element(slide)
    
    def create_section_slide(self, section_title: str, description: str = "") -> None:
        """Create a section divider slide"""
        slide_layout = self.prs.slide_layouts[2]  # Section header layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set section title
        title_shape = slide.shapes.title
        title_shape.text = section_title
        self.apply_font_style(title_shape.text_frame, "title", self.color_scheme["primary"])
        
        # Add description if provided
        if description and len(slide.placeholders) > 1:
            desc_shape = slide.placeholders[1]
            desc_shape.text = description
            self.apply_font_style(desc_shape.text_frame, "subtitle", self.color_scheme["secondary"])
        
        # Add background accent
        self._add_background_accent(slide)
    
    def create_comparison_slide(self, title: str, left_title: str, left_content: List[str], 
                               right_title: str, right_content: List[str]) -> None:
        """Create a two-column comparison slide"""
        slide_layout = self.prs.slide_layouts[3]  # Two content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set main title
        title_shape = slide.shapes.title
        title_shape.text = title
        self.apply_font_style(title_shape.text_frame, "heading", self.color_scheme["primary"])
        
        # Left column
        left_shape = slide.placeholders[1]
        left_text = f"{left_title}\n\n" + "\n".join([f"• {item}" for item in left_content])
        left_shape.text = left_text
        self.apply_font_style(left_shape.text_frame, "body")
        
        # Right column
        right_shape = slide.placeholders[2]
        right_text = f"{right_title}\n\n" + "\n".join([f"• {item}" for item in right_content])
        right_shape.text = right_text
        self.apply_font_style(right_shape.text_frame, "body")
    
    def create_conclusion_slide(self, title: str = "Key Takeaways", 
                              takeaways: List[str] = None) -> None:
        """Create a conclusion slide with key takeaways"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        self.apply_font_style(title_shape.text_frame, "heading", self.color_scheme["primary"])
        
        # Set takeaways
        if takeaways:
            content_shape = slide.placeholders[1]
            content_text = "\n".join([f"✓ {item}" for item in takeaways])
            content_shape.text = content_text
            self.apply_font_style(content_shape.text_frame, "body")
        
        # Add special styling for conclusion
        self._add_conclusion_styling(slide)
    
    def create_thank_you_slide(self, main_text: str = "Thank You", 
                              sub_text: str = "Questions?") -> None:
        """Create a professional thank you slide"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set main text
        title_shape = slide.shapes.title
        title_shape.text = main_text
        self.apply_font_style(title_shape.text_frame, "title", self.color_scheme["primary"])
        
        # Set sub text
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = sub_text
            self.apply_font_style(subtitle_shape.text_frame, "subtitle", self.color_scheme["secondary"])
        
        # Add decorative elements
        self._add_thank_you_styling(slide)
    
    def create_image_slide(self, title: str, image_path: str, caption: str = "", 
                          layout_style: str = "center") -> None:
        """
        Create a slide with an image as the main content
        
        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption for the image
            layout_style: Layout style ('center', 'left', 'right', 'full')
        """
        try:
            # Use blank layout but remove any existing placeholders
            slide_layout = self.prs.slide_layouts[6]  # Completely blank layout
            slide = self.prs.slides.add_slide(slide_layout)
            
            # Clear any existing placeholders to avoid "Click to add title" issues
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        slide.shapes.element.remove(shape.element)
                    except:
                        pass
            
            # Add title with proper styling
            title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2),  # position (slightly higher)
                Inches(9), Inches(0.8)     # size (slightly smaller)
            )
            title_shape.text = title
            self.apply_font_style(title_shape.text_frame, "heading", self.color_scheme["primary"])
            
            # Ensure title is left-aligned and visible
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Check if image file exists
            if not os.path.exists(image_path):
                logger.warning(f"Image file not found: {image_path}")
                # Create a placeholder text instead
                placeholder_shape = slide.shapes.add_textbox(
                    Inches(2), Inches(3),
                    Inches(6), Inches(2)
                )
                placeholder_shape.text = f"Image not found:\n{Path(image_path).name}"
                self.apply_font_style(placeholder_shape.text_frame, "body", self.color_scheme["secondary"])
                return
            
            # Add image based on layout style
            if layout_style == "center":
                # Center the image
                img_shape = slide.shapes.add_picture(
                    image_path,
                    Inches(2), Inches(1.5),    # position
                    width=Inches(6)            # width (height will be proportional)
                )
            elif layout_style == "left":
                # Image on left, space for text on right
                img_shape = slide.shapes.add_picture(
                    image_path,
                    Inches(0.5), Inches(1.5),
                    width=Inches(4.5)
                )
            elif layout_style == "right":
                # Image on right, space for text on left
                img_shape = slide.shapes.add_picture(
                    image_path,
                    Inches(5), Inches(1.5),
                    width=Inches(4.5)
                )
            elif layout_style == "full":
                # Full slide image with overlay title
                img_shape = slide.shapes.add_picture(
                    image_path,
                    Inches(0), Inches(1.2),
                    width=Inches(10)
                )
            else:
                # Default to center
                img_shape = slide.shapes.add_picture(
                    image_path,
                    Inches(2), Inches(1.5),
                    width=Inches(6)
                )
            
            # Add caption if provided
            if caption:
                if layout_style == "full":
                    # Caption overlay for full image
                    caption_shape = slide.shapes.add_textbox(
                        Inches(0.5), Inches(6.5),
                        Inches(9), Inches(0.8)
                    )
                else:
                    # Caption below image
                    caption_shape = slide.shapes.add_textbox(
                        Inches(1), Inches(6),
                        Inches(8), Inches(0.8)
                    )
                
                caption_shape.text = caption
                self.apply_font_style(caption_shape.text_frame, "caption", self.color_scheme["secondary"])
                
                # Center align caption
                for paragraph in caption_shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
            
            # Add accent element
            self._add_accent_element(slide)
            
        except Exception as e:
            logger.error(f"Error creating image slide: {str(e)}")
            # Fall back to creating a text slide with error message
            self.create_content_slide(title, [f"Error loading image: {Path(image_path).name}"])
    
    def create_image_content_slide(self, title: str, image_path: str, content: List[str], 
                                  image_position: str = "right") -> None:
        """
        Create a slide with both image and bullet point content
        
        Args:
            title: Slide title
            image_path: Path to the image file
            content: List of bullet points
            image_position: Position of image ('left' or 'right')
        """
        try:
            # Use a layout that has title and content placeholders, then customize
            slide_layout = self.prs.slide_layouts[6]  # Try blank layout without placeholders
            slide = self.prs.slides.add_slide(slide_layout)
            
            # Clear any existing placeholders to avoid "Click to add title" issues
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        slide.shapes.element.remove(shape.element)
                    except:
                        pass
            
            # Add title with better positioning and styling
            title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2),  # Slightly higher position
                Inches(9), Inches(0.8)    # Slightly smaller height
            )
            title_shape.text = title
            self.apply_font_style(title_shape.text_frame, "heading", self.color_scheme["primary"])
            
            # Center align the title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
            
            if image_position == "right":
                # Content on left, image on right with proper spacing
                content_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.2),   # Start position
                    Inches(4.5), Inches(5.5)    # Smaller width to prevent cutoff
                )
                
                if os.path.exists(image_path):
                    img_shape = slide.shapes.add_picture(
                        image_path,
                        Inches(5.2), Inches(1.2),   # More space from content
                        width=Inches(4.5)           # Larger image
                    )
                else:
                    # Image placeholder
                    img_placeholder = slide.shapes.add_textbox(
                        Inches(5.5), Inches(3),
                        Inches(4), Inches(1.5)
                    )
                    img_placeholder.text = f"Image:\n{Path(image_path).name}"
                    self.apply_font_style(img_placeholder.text_frame, "caption", self.color_scheme["secondary"])
                    
            else:  # image_position == "left"
                # Image on left, content on right with proper spacing
                if os.path.exists(image_path):
                    img_shape = slide.shapes.add_picture(
                        image_path,
                        Inches(0.5), Inches(1.2),   # Start position
                        width=Inches(4.5)           # Larger image
                    )
                else:
                    # Image placeholder
                    img_placeholder = slide.shapes.add_textbox(
                        Inches(0.5), Inches(3),
                        Inches(4.5), Inches(1.5)
                    )
                    img_placeholder.text = f"Image:\n{Path(image_path).name}"
                    self.apply_font_style(img_placeholder.text_frame, "caption", self.color_scheme["secondary"])
                
                content_shape = slide.shapes.add_textbox(
                    Inches(5.2), Inches(1.2),    # More space from image
                    Inches(4.5), Inches(5.5)     # Prevent right-side cutoff
                )
            
            # Add bullet point content with proper formatting
            content_text = "\n\n".join([f"• {item}" for item in content])  # Double spacing between bullets
            content_shape.text = content_text
            self.apply_font_style(content_shape.text_frame, "body")
            
            # Enable word wrapping and improve text layout
            content_shape.text_frame.word_wrap = True
            content_shape.text_frame.auto_size = None  # Disable auto-sizing to prevent overflow
            
            # Set margins for better text spacing
            content_shape.text_frame.margin_left = Inches(0.15)
            content_shape.text_frame.margin_right = Inches(0.15)
            content_shape.text_frame.margin_top = Inches(0.2)
            content_shape.text_frame.margin_bottom = Inches(0.2)
            
            # Improve paragraph spacing for readability
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.space_after = Pt(8)  # Space after each paragraph
            
            # Add accent element
            self._add_accent_element(slide)
            
        except Exception as e:
            logger.error(f"Error creating image content slide: {str(e)}")
            # Fall back to regular content slide
            self.create_content_slide(title, content)
    
    def create_image_comparison_slide(self, title: str, left_image: str, left_title: str, 
                                    right_image: str, right_title: str, 
                                    left_caption: str = "", right_caption: str = "") -> None:
        """
        Create a slide comparing two images side by side
        
        Args:
            title: Main slide title
            left_image: Path to left image
            left_title: Title for left image
            right_image: Path to right image  
            right_title: Title for right image
            left_caption: Optional caption for left image
            right_caption: Optional caption for right image
        """
        try:
            slide_layout = self.prs.slide_layouts[5]  # Blank layout
            slide = self.prs.slides.add_slide(slide_layout)
            
            # Add main title
            title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(9), Inches(1)
            )
            title_shape.text = title
            self.apply_font_style(title_shape.text_frame, "heading", self.color_scheme["primary"])
            
            # Left side
            left_title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3),
                Inches(4.5), Inches(0.5)
            )
            left_title_shape.text = left_title
            self.apply_font_style(left_title_shape.text_frame, "body", self.color_scheme["primary"])
            left_title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            if os.path.exists(left_image):
                left_img = slide.shapes.add_picture(
                    left_image,
                    Inches(0.5), Inches(2),
                    width=Inches(4.5)
                )
            else:
                left_placeholder = slide.shapes.add_textbox(
                    Inches(0.5), Inches(3),
                    Inches(4.5), Inches(2)
                )
                left_placeholder.text = f"Image not found:\n{Path(left_image).name}"
                self.apply_font_style(left_placeholder.text_frame, "caption", self.color_scheme["secondary"])
            
            if left_caption:
                left_caption_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(5.5),
                    Inches(4.5), Inches(0.8)
                )
                left_caption_shape.text = left_caption
                self.apply_font_style(left_caption_shape.text_frame, "caption", self.color_scheme["secondary"])
                left_caption_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Right side
            right_title_shape = slide.shapes.add_textbox(
                Inches(5), Inches(1.3),
                Inches(4.5), Inches(0.5)
            )
            right_title_shape.text = right_title
            self.apply_font_style(right_title_shape.text_frame, "body", self.color_scheme["primary"])
            right_title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            if os.path.exists(right_image):
                right_img = slide.shapes.add_picture(
                    right_image,
                    Inches(5), Inches(2),
                    width=Inches(4.5)
                )
            else:
                right_placeholder = slide.shapes.add_textbox(
                    Inches(5), Inches(3),
                    Inches(4.5), Inches(2)
                )
                right_placeholder.text = f"Image not found:\n{Path(right_image).name}"
                self.apply_font_style(right_placeholder.text_frame, "caption", self.color_scheme["secondary"])
            
            if right_caption:
                right_caption_shape = slide.shapes.add_textbox(
                    Inches(5), Inches(5.5),
                    Inches(4.5), Inches(0.8)
                )
                right_caption_shape.text = right_caption
                self.apply_font_style(right_caption_shape.text_frame, "caption", self.color_scheme["secondary"])
                right_caption_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add dividing line
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.8), Inches(1.5),
                Inches(0.05), Inches(4.5)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = self.color_scheme["accent"]
            divider.line.fill.background()
            
        except Exception as e:
            logger.error(f"Error creating image comparison slide: {str(e)}")
            # Fall back to text comparison
            self.create_comparison_slide(title, left_title, [f"Image: {Path(left_image).name}"], 
                                       right_title, [f"Image: {Path(right_image).name}"])
    
    def _add_accent_line(self, slide):
        """Add a decorative accent line to the slide"""
        try:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(2.8),  # position
                Inches(9), Inches(0.1)      # size
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.color_scheme["accent"]
            line.line.fill.background()
        except:
            pass  # Skip if layout doesn't support shapes
    
    def _add_accent_element(self, slide):
        """Add a small accent element to content slides"""
        try:
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.3), Inches(1.5),   # position
                Inches(0.2), Inches(1.5)    # size
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self.color_scheme["accent"]
            accent.line.fill.background()
        except:
            pass
    
    def _add_background_accent(self, slide):
        """Add background accent for section slides"""
        try:
            bg_accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),        # position
                Inches(10), Inches(0.5)      # size
            )
            bg_accent.fill.solid()
            bg_accent.fill.fore_color.rgb = self.color_scheme["primary"]
            bg_accent.line.fill.background()
        except:
            pass
    
    def _add_conclusion_styling(self, slide):
        """Add special styling for conclusion slides"""
        try:
            # Add a subtle background accent
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.2), Inches(2.5),
                Inches(0.1), Inches(4)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self.color_scheme["accent"]
            accent.line.fill.background()
        except:
            pass
    
    def _add_thank_you_styling(self, slide):
        """Add decorative elements to thank you slide"""
        try:
            # Add corner accents
            for pos in [(0.2, 0.2), (9.5, 0.2), (0.2, 7), (9.5, 7)]:
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(pos[0]), Inches(pos[1]),
                    Inches(0.3), Inches(0.3)
                )
                accent.fill.solid()
                accent.fill.fore_color.rgb = self.color_scheme["accent"]
                accent.line.fill.background()
        except:
            pass
    
    def save(self, filename: str):
        """Save the presentation"""
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        # Ensure proper path handling
        from pathlib import Path
        file_path = Path(filename)
        self.prs.save(str(file_path))

class BusinessTemplate(PPTTemplate):
    """Professional business presentation template"""
    
    def __init__(self):
        super().__init__("business_blue")

class TechTemplate(PPTTemplate):
    """Modern technology presentation template"""
    
    def __init__(self):
        super().__init__("tech_green")

class CreativeTemplate(PPTTemplate):
    """Creative/startup presentation template"""
    
    def __init__(self):
        super().__init__("modern_purple")

class CorporateTemplate(PPTTemplate):
    """Corporate presentation template"""
    
    def __init__(self):
        super().__init__("corporate_red")

# Template factory function
def create_template(template_type: str = "business") -> PPTTemplate:
    """
    Create a presentation template of the specified type
    
    Args:
        template_type: Type of template ("business", "tech", "creative", "corporate")
        
    Returns:
        PPTTemplate instance
    """
    templates = {
        "business": BusinessTemplate,
        "tech": TechTemplate,
        "creative": CreativeTemplate,
        "corporate": CorporateTemplate
    }
    
    template_class = templates.get(template_type, BusinessTemplate)
    return template_class()

# Example usage function
def create_sample_presentation():
    """Create a sample presentation demonstrating the template system"""
    
    # Create a business template
    ppt = create_template("business")
    
    # Title slide
    ppt.create_title_slide(
        "Sample Business Presentation",
        "Demonstrating Professional Template System"
    )
    
    # Section slide
    ppt.create_section_slide(
        "Introduction",
        "Setting the stage for our presentation"
    )
    
    # Content slide with bullet points
    ppt.create_content_slide(
        "Key Features",
        [
            "Professional design templates",
            "Consistent color schemes",
            "Modern typography",
            "Flexible layout options",
            "Easy customization"
        ]
    )
    
    # Comparison slide
    ppt.create_comparison_slide(
        "Before vs After",
        "Before Template System",
        ["Inconsistent designs", "Manual formatting", "Time-consuming"],
        "After Template System", 
        ["Professional appearance", "Automated styling", "Quick generation"]
    )
    
    # Conclusion slide
    ppt.create_conclusion_slide(
        "Key Benefits",
        [
            "Saves time and effort",
            "Ensures professional quality",
            "Maintains brand consistency",
            "Easy to use and customize"
        ]
    )
    
    # Thank you slide
    ppt.create_thank_you_slide()
    
    # Save the presentation
    ppt.save("sample_business_presentation.pptx")
    print("Sample presentation created: sample_business_presentation.pptx")

if __name__ == "__main__":
    create_sample_presentation()

